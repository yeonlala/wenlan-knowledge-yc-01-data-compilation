"""
Mock 文件占位内容：生成可被 Word / Excel / PDF 阅读器打开的最小合法文件（非 0 字节）。

正文优先来自包内目录 mock_templates/rXX.md（或 .txt），便于维护「正规范例」；
若文件不存在则使用内置短占位。中文 PDF 使用 ReportLab + 系统中文字体。
"""

from __future__ import annotations

import io
import os
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional

MOCK_TEMPLATES_DIR = Path(__file__).resolve().parent / "mock_templates"


def template_source_file(rid: str) -> Optional[Path]:
    """返回仓库内可编辑范例文件路径（.md 优先于 .txt）。"""
    for ext in (".md", ".txt"):
        p = MOCK_TEMPLATES_DIR / f"{rid}{ext}"
        if p.is_file():
            return p
    return None


def load_rule_template(rid: str) -> Optional[str]:
    p = template_source_file(rid)
    if p is None:
        return None
    return p.read_text(encoding="utf-8")


def xml_escape(s: str) -> str:
    return (
        s.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def find_cjk_font_path() -> Optional[str]:
    """优先 Windows 微软雅黑 / 黑体 / 宋体；Linux 常见 Noto / 文泉驿。"""
    candidates: List[str] = []
    if os.name == "nt":
        windir = os.environ.get("WINDIR", r"C:\Windows")
        fonts = os.path.join(windir, "Fonts")
        candidates.extend(
            [
                os.path.join(fonts, "msyh.ttc"),
                os.path.join(fonts, "msyhbd.ttc"),
                os.path.join(fonts, "simhei.ttf"),
                os.path.join(fonts, "simsun.ttc"),
                os.path.join(fonts, "simsunb.ttf"),
            ]
        )
    candidates.extend(
        [
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        ]
    )
    for p in candidates:
        if os.path.isfile(p):
            return p
    return None


def _footer_block(rid: str, note: str, customer: str, project: str) -> str:
    return (
        "\n\n---\n"
        "【以下为脚本追加页脚，可整段删除】\n"
        f"规则编号：{rid}\n"
        f"材料说明：{note}\n"
        f"客户单位：{customer}\n"
        f"项目名称：{project}\n"
    )


def compose_main_body(
    rule: Dict[str, Any], customer: str, project: str
) -> str:
    """合并范例模板与页脚；无模板时用内置短正文。"""
    rid = str(rule["id"])
    note = str(rule.get("note", ""))
    tmpl = load_rule_template(rid)
    foot = _footer_block(rid, note, customer, project)
    if tmpl is not None:
        return tmpl.strip() + foot
    return (
        "（未找到 mock_templates/ 下对应范例文件，此为内置短占位。）\n"
        f"请将正式稿放入 tobacco_kb/mock_templates/{rid}.md 后重新生成。\n\n"
        f"规则编号：{rid}\n"
        f"材料说明：{note}\n"
        f"客户单位：{customer}\n"
        f"项目名称：{project}"
    )


def docx_bytes_from_paragraphs(paragraphs: List[str]) -> bytes:
    """多段 Word，每段一个段落。"""
    parts: List[str] = []
    for para in paragraphs:
        esc = xml_escape(para)
        parts.append(
            f'<w:p><w:r><w:t xml:space="preserve">{esc}</w:t></w:r></w:p>'
        )
    inner = "".join(parts)
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{inner}</w:body></w:document>"
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" '
        'ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.'
        'wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/'
        'officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
        "</Relationships>"
    )
    word_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types.encode("utf-8"))
        zf.writestr("_rels/.rels", rels.encode("utf-8"))
        zf.writestr("word/_rels/document.xml.rels", word_rels.encode("utf-8"))
        zf.writestr("word/document.xml", document.encode("utf-8"))
    return buf.getvalue()


def minimal_docx_bytes(title: str, body: str) -> bytes:
    """标题 + 正文；正文按空行分段，否则按单行分段。"""
    paras: List[str] = [title]
    b = body.strip()
    if "\n\n" in b:
        for p in b.split("\n\n"):
            s = p.strip()
            if s:
                paras.append(s)
    else:
        for ln in b.split("\n"):
            paras.append(ln.rstrip())
    return docx_bytes_from_paragraphs(paras)


def docx_pretty_or_legacy(title: str, body: str) -> bytes:
    """优先将 Markdown 渲染为带样式的 Word；缺少 python-docx 时退回简易 OOXML。"""
    try:
        from .markdown_to_docx import has_python_docx, render_markdown_to_docx_bytes

        if has_python_docx():
            return render_markdown_to_docx_bytes(title, body)
    except Exception:
        pass
    return minimal_docx_bytes(title, body)


def minimal_pdf_bytes_unicode(full_text: str) -> bytes:
    """
    中文 PDF：ReportLab + TrueType / TTC；支持多页与长文换行。
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen import canvas

    font_path = find_cjk_font_path()
    if not font_path:
        raise OSError(
            "未找到系统中文字体文件（如 Windows 下的 msyh.ttc）。"
            "无法生成可读的中文 PDF，请安装字体或改用：--ext .docx"
        )

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)
    font_name = "MockCJK"
    if font_path.lower().endswith(".ttc"):
        pdfmetrics.registerFont(TTFont(font_name, font_path, subfontIndex=0))
    else:
        pdfmetrics.registerFont(TTFont(font_name, font_path))

    w, h = A4
    c.setFont(font_name, 11)
    y = h - 72
    bottom = 56
    line_h = 14
    max_chars = 26

    full_text = full_text.replace("\r\n", "\n").replace("\r", "\n")
    for raw in full_text.split("\n"):
        line = raw.rstrip()
        if not line:
            y -= line_h // 2
            if y < bottom:
                c.showPage()
                c.setFont(font_name, 11)
                y = h - 72
            continue
        while len(line) > max_chars:
            chunk = line[:max_chars]
            line = line[max_chars:]
            if y < bottom:
                c.showPage()
                c.setFont(font_name, 11)
                y = h - 72
            c.drawString(72, y, chunk)
            y -= line_h
        if y < bottom:
            c.showPage()
            c.setFont(font_name, 11)
            y = h - 72
        c.drawString(72, y, line)
        y -= line_h

    c.save()
    return buf.getvalue()


def minimal_pdf_bytes(full_text: str) -> bytes:
    try:
        return minimal_pdf_bytes_unicode(full_text)
    except ImportError as e:
        raise ImportError(
            "生成中文 PDF 需要先安装：pip install reportlab\n"
            "安装后仍会自动使用系统中的微软雅黑等字体。\n"
            "若暂不安装，请使用默认：--ext .docx（Word 可完整显示中文）。"
        ) from e


def pdf_pretty_or_legacy(title: str, body: str) -> bytes:
    """优先将 Markdown 渲染为带样式的 PDF（与 docx 共用解析）；失败则退回按行排版。"""
    try:
        from .markdown_to_pdf import has_markdown_pdf, render_markdown_to_pdf_bytes

        if has_markdown_pdf():
            return render_markdown_to_pdf_bytes(title, body)
    except Exception:
        pass
    full_pdf = f"{title}\n\n{body}"
    return minimal_pdf_bytes(full_pdf)


def minimal_pptx_bytes(title: str, body: str) -> bytes:
    try:
        from io import BytesIO

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        try:
            layout = prs.slide_layouts[6]
        except IndexError:
            layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
        tf = box.text_frame
        tf.text = title
        p = tf.add_paragraph()
        p.text = body[:12000]
        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        return docx_pretty_or_legacy(
            title + "（未安装 python-pptx，以下为 Word 占位内容）",
            body,
        )


def minimal_xlsx_bytes(sheet_title: str, cell_a1: str, cell_a2: str) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_title[:31] if sheet_title else "占位"
    ws["A1"] = cell_a1
    ws["A2"] = cell_a2
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def minimal_xlsx_column(sheet_title: str, column_lines: List[str]) -> bytes:
    """A 列自上而下写入多行（范例正文逐行）。"""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_title[:31] if sheet_title else "占位"
    for i, line in enumerate(column_lines, start=1):
        ws.cell(row=i, column=1, value=line)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def minimal_zip_bytes(readme_text: str) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("readme.txt", readme_text.encode("utf-8"))
    return buf.getvalue()


def bytes_for_mock_file(
    suffix: str,
    rule: Dict[str, Any],
    customer: str,
    project: str,
) -> bytes:
    """
    根据扩展名返回占位二进制。
    rule: PHASE1_RULES 中一条 dict。
    """
    ext = suffix.lower()
    if not ext.startswith("."):
        ext = "." + ext

    rid = str(rule["id"])
    label = str(rule.get("label", rid))
    if rid == "aux":
        line_title = f"【命名测试】{label}"
        main_body = default_body_for_aux(customer, project, label)
    else:
        line_title = f"【模拟资料】{label}"
        main_body = compose_main_body(rule, customer, project)

    if ext == ".pdf":
        return pdf_pretty_or_legacy(line_title, main_body)

    if ext in (".docx", ".doc"):
        return docx_pretty_or_legacy(line_title, main_body)

    if ext in (".xlsx", ".xls"):
        try:
            import openpyxl  # noqa: F401
        except ImportError:
            return docx_pretty_or_legacy(
                line_title + "（无法生成 xlsx 时的 Word 占位）",
                main_body + "\n\n提示：安装 openpyxl 后可生成 Excel。",
            )
        lines = main_body.splitlines()
        return minimal_xlsx_column("范例内容", lines)

    if ext in (".pptx", ".ppt"):
        return minimal_pptx_bytes(line_title, main_body)

    if ext in (".png", ".jpg", ".jpeg"):
        if ext == ".png":
            return (
                b"\x89PNG\r\n\x1a\n"
                b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
                b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89"
                b"\x00\x00\x00\nIDATx\x9cc\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4"
                b"\x00\x00\x00\x00IEND\xaeB`\x82"
            )
        return (
            b"\xff\xd8\xff\xe0\x00\x10JFIF\x00\x01\x01\x00\x00\x01\x00\x01\x00\x00"
            b"\xff\xdb\x00C\x00\x08\x06\x06\x07\x06\x05\x08\x07\x07\x07\t\t\x08\n"
            b"\x0c\x14\r\x0c\x0b\x0b\x0c\x19\x12\x13\x0f\x14\x1d\x1a\x1f\x1e\x1d\x1a"
            b"\x1c\x1c $.\' \",#\x1c\x1c(7),01444\x1f\'9=82<.342\xff\xc0\x00\x11\x08"
            b"\x00\x01\x00\x01\x01\x01\x11\x00\x02\x11\x01\x03\x11\x01\xff\xc4\x00\x14"
            b"\x00\x01\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00"
            b"\x00\x08\xff\xc4\x00\x14\x10\x01\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00"
            b"\x00\x00\x00\x00\x00\x00\x00\x00\xff\xda\x00\x0c\x03\x01\x00\x02\x11\x03"
            b"\x11\x00\x3f\x00\xaa\xff\xd9"
        )

    if ext == ".zip":
        return minimal_zip_bytes(f"{line_title}\n\n{main_body}\n")

    if ext == ".rar":
        return (
            "此为模拟压缩包占位说明（扩展名为 rar 时建议改用 zip 测试）。\n"
            f"{line_title}\n\n{main_body}"
        ).encode("utf-8")

    return (
        f"{line_title}\n\n{main_body}\n\n"
        f"（扩展名为 {ext}，本文件为 UTF-8 文本占位。）"
    ).encode("utf-8")


def default_body_for_aux(customer: str, project: str, label: str) -> str:
    """bad_names 等无 PHASE1 规则时的正文。"""
    return (
        f"标签：{label}\n"
        f"客户单位：{customer}\n"
        f"项目名称：{project}\n"
        "本文件为命名规则测试用占位内容。"
    )
